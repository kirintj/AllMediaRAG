"""PPTX 文件读取器"""
from __future__ import annotations


class PptxReader:
    """读取 .pptx 文件，提取每页文本框 + 表格 + 备注"""

    def read(self, file_path: str) -> str:
        from pptx import Presentation
        prs = Presentation(file_path)
        parts = []
        for i, slide in enumerate(prs.slides, 1):
            texts = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            texts.append(text)
                if shape.has_table:
                    table = shape.table
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        texts.append(" | ".join(cells))
            if slide.has_notes_slide:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    texts.append(f"[备注] {notes}")
            if texts:
                parts.append(f"## 第{i}页\n" + "\n".join(texts))
        return "\n\n".join(parts)

    def supported_extensions(self) -> list[str]:
        return [".pptx"]
